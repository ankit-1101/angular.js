from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE


def set_text_style(text_frame, font_size=22, bold=False, color=RGBColor(40, 40, 40)):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color


def add_title_slide(prs: Presentation, title: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    # style
    set_text_style(slide.shapes.title.text_frame, font_size=44, bold=True, color=RGBColor(20, 40, 90))
    set_text_style(slide.placeholders[1].text_frame, font_size=20, bold=False, color=RGBColor(70, 70, 70))


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1]
    tf = body.text_frame
    tf.clear()
    if not bullets:
        return slide
    # First bullet
    p = tf.paragraphs[0]
    p.text = bullets[0]
    # Subsequent bullets
    for b in bullets[1:]:
        p = tf.add_paragraph()
        p.text = b
        p.level = 0
    # style
    set_text_style(slide.shapes.title.text_frame, font_size=32, bold=True, color=RGBColor(20, 40, 90))
    set_text_style(tf, font_size=22, bold=False, color=RGBColor(40, 40, 40))
    return slide


def add_two_content_slide(prs: Presentation, title: str, left_bullets: list[str], right_bullets: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[3])  # Two Content
    slide.shapes.title.text = title

    left = slide.placeholders[1]
    right = slide.placeholders[2]

    left_tf = left.text_frame
    left_tf.clear()
    if left_bullets:
        left_tf.paragraphs[0].text = left_bullets[0]
        for b in left_bullets[1:]:
            p = left_tf.add_paragraph()
            p.text = b

    right_tf = right.text_frame
    right_tf.clear()
    if right_bullets:
        right_tf.paragraphs[0].text = right_bullets[0]
        for b in right_bullets[1:]:
            p = right_tf.add_paragraph()
            p.text = b

    set_text_style(slide.shapes.title.text_frame, font_size=32, bold=True, color=RGBColor(20, 40, 90))
    set_text_style(left_tf, font_size=22, bold=False, color=RGBColor(40, 40, 40))
    set_text_style(right_tf, font_size=22, bold=False, color=RGBColor(40, 40, 40))
    return slide


def add_simple_anatomy_diagram(slide, left=Inches(6.3), top=Inches(1.6)):
    # Simple schematic: circle (humeral head), rectangle (acromion), line (supraspinatus)
    shapes = slide.shapes
    # Humeral head
    humeral_head = shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left, top, Inches(1.3), Inches(1.3))
    humeral_head.fill.solid()
    humeral_head.fill.fore_color.rgb = RGBColor(230, 240, 255)
    humeral_head.line.color.rgb = RGBColor(20, 40, 90)
    humeral_head.text = "Humeral\nHead"
    set_text_style(humeral_head.text_frame, font_size=12, bold=False, color=RGBColor(20, 40, 90))

    # Acromion
    acromion = shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left - Inches(1.8), top - Inches(0.7), Inches(2.0), Inches(0.5))
    acromion.fill.solid()
    acromion.fill.fore_color.rgb = RGBColor(245, 245, 245)
    acromion.line.color.rgb = RGBColor(120, 120, 120)
    acromion.text = "Acromion"
    set_text_style(acromion.text_frame, font_size=12, bold=False, color=RGBColor(80, 80, 80))

    # Supraspinatus tendon (simple narrow rectangle)
    supra = shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left - Inches(0.6), top + Inches(0.45), Inches(0.8), Inches(0.15))
    supra.fill.solid()
    supra.fill.fore_color.rgb = RGBColor(200, 30, 30)
    supra.line.color.rgb = RGBColor(150, 0, 0)
    supra.text = "Supraspinatus\nTendon"
    set_text_style(supra.text_frame, font_size=10, bold=False, color=RGBColor(255, 255, 255))


def add_simple_timeline(slide):
    shapes = slide.shapes
    left = Inches(0.6)
    top = Inches(4.8)
    width = Inches(8.8)
    height = Inches(0.25)
    bar = shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(220, 226, 240)
    bar.line.color.rgb = RGBColor(20, 40, 90)

    milestones = [
        ("0-2 wks", 0.0),
        ("2-6 wks", 0.25),
        ("6-12 wks", 0.55),
        ("12+ wks", 0.85),
    ]
    for label, rel_x in milestones:
        x = left + width * rel_x
        m = shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, x - Inches(0.15), top - Inches(0.35), Inches(0.3), Inches(0.3))
        m.fill.solid()
        m.fill.fore_color.rgb = RGBColor(20, 40, 90)
        m.line.color.rgb = RGBColor(20, 40, 90)
        tbox = shapes.add_textbox(x - Inches(0.5), top - Inches(0.8), Inches(1.2), Inches(0.4))
        tbox.text_frame.text = label
        set_text_style(tbox.text_frame, font_size=12, bold=True, color=RGBColor(20, 40, 90))


def build_presentation(output_path: str):
    prs = Presentation()

    # Title slide
    add_title_slide(
        prs,
        title="Rotator Cuff Injury: Surgical and Physiotherapy Management",
        subtitle="Final-year BPT Presentation\nName: ____________  |  College: ____________  |  Date: ____________",
    )

    # Introduction
    add_bullet_slide(
        prs,
        title="Introduction",
        bullets=[
            "Rotator cuff = group of 4 muscles stabilizing the glenohumeral joint",
            "Injuries: tendinopathy, partial- or full-thickness tears",
            "Common in overhead activities and with aging",
            "Results in pain, weakness, and functional limitation",
        ],
    )

    # Anatomy & function (with simple diagram)
    slide = add_bullet_slide(
        prs,
        title="Anatomy & Function (Quick Recap)",
        bullets=[
            "SITS: Supraspinatus, Infraspinatus, Teres minor, Subscapularis",
            "Primary roles: humeral head centering, shoulder elevation and rotation",
            "Subacromial space relevance in impingement",
        ],
    )
    add_simple_anatomy_diagram(slide)

    # Causes & risk factors
    add_bullet_slide(
        prs,
        title="Causes & Risk Factors",
        bullets=[
            "Degenerative: age-related tendon wear, hypovascular zones",
            "Overuse: repetitive overhead work/sports (e.g., swimming, tennis)",
            "Traumatic: fall on outstretched hand, sudden eccentric load",
            "Risk factors: smoking, diabetes, hyperlipidemia, poor posture",
        ],
    )

    # Symptoms & clinical presentation
    add_bullet_slide(
        prs,
        title="Symptoms & Clinical Presentation",
        bullets=[
            "Pain on lateral shoulder/upper arm, worse with overhead activity",
            "Night pain, difficulty sleeping on affected side",
            "Weakness in abduction/external rotation",
            "Painful arc, crepitus; positive impingement signs",
        ],
    )

    # Assessment (tests & imaging)
    add_bullet_slide(
        prs,
        title="Assessment: Tests & Imaging",
        bullets=[
            "Special tests: Jobe (Empty Can), Hawkins–Kennedy, Neer, External Rotation Lag",
            "Palpation tenderness over greater tuberosity and subacromial region",
            "ROM & strength (MMT) with scapular control assessment",
            "Imaging: X-ray (acromial spur), Ultrasound, MRI (tear size/retreat)",
        ],
    )

    # Surgical management: indications
    add_bullet_slide(
        prs,
        title="Surgical Management: Indications",
        bullets=[
            "Persistent pain/functional loss after ≥3–6 months of structured rehab",
            "Acute traumatic full-thickness tear in active individuals",
            "Large tears with retraction, failed previous repair",
            "Associated lesions (biceps pathology, impingement) requiring intervention",
        ],
    )

    # Surgical procedures
    add_two_content_slide(
        prs,
        title="Surgical Procedures (Overview)",
        left_bullets=[
            "Arthroscopic debridement/subacromial decompression",
            "Arthroscopic or mini-open rotator cuff repair",
            "Tendon mobilization, footprint preparation, suture anchors",
            "Biceps procedures: tenotomy/tenodesis when indicated",
        ],
        right_bullets=[
            "Small: <1 cm; Medium: 1–3 cm; Large: 3–5 cm; Massive: >5 cm",
            "Single-row vs. double-row configurations",
            "Concomitant acromioplasty for impingement morphology",
            "Reverse TSA for cuff tear arthropathy (select cases)",
        ],
    )

    # Surgical outcomes/complications
    add_bullet_slide(
        prs,
        title="Surgery: Outcomes & Considerations",
        bullets=[
            "Pain relief is common; strength recovery depends on tear size/quality",
            "Re-tear risk: age, tendon quality, large/massive tears",
            "Complications: stiffness, infection, anchor issues, biceps pain",
            "Adherence to post-op protocol is crucial for tendon healing",
        ],
    )

    # Physiotherapy management: overview
    add_bullet_slide(
        prs,
        title="Physiotherapy Management: Overview",
        bullets=[
            "Tailor program to diagnosis: non-operative vs. post-operative",
            "Stage-based progression guided by pain, ROM, and healing timelines",
            "Key pillars: pain control, mobility, motor control, strength, function",
            "Education: activity modification and load management",
        ],
    )

    # Pre-op physiotherapy
    add_bullet_slide(
        prs,
        title="Pre-operative Physiotherapy (""Prehab"" )",
        bullets=[
            "Goals: pain modulation, maintain ROM, optimize scapular mechanics",
            "Interventions: pendulums, AAROM within tolerance, postural cues",
            "Isometrics for deltoid/RC within pain limits; scapular setting",
            "Education: sling use post-op, sleep positioning, expectations",
        ],
    )

    # Post-op phases (early)
    slide = add_bullet_slide(
        prs,
        title="Post-op Phase I (0–2 weeks)",
        bullets=[
            "Protection: sling/abduction pillow as prescribed",
            "Pain/edema control: cold therapy, supported positioning",
            "Passive ROM only (per protocol): flexion, ER in scapular plane",
            "Pendulums, distal AROM (elbow/wrist/hand)",
        ],
    )
    add_simple_timeline(slide)

    slide = add_bullet_slide(
        prs,
        title="Post-op Phase II (2–6 weeks)",
        bullets=[
            "Progress PROM; initiate AAROM (pulleys, wand) within limits",
            "Scapular stability: retraction/depression drills, closed-chain weight shifts",
            "Avoid active shoulder elevation if repair not ready",
            "Monitor for stiffness; respect pain and night symptoms",
        ],
    )
    add_simple_timeline(slide)

    slide = add_bullet_slide(
        prs,
        title="Post-op Phase III (6–12 weeks)",
        bullets=[
            "Transition to AROM then light strengthening as cleared",
            "Isometrics → bands for ER/IR, rows; avoid painful arcs",
            "Neuromuscular control: scapulohumeral rhythm, serratus/low trap",
            "Functional reaching patterns below shoulder height initially",
        ],
    )
    add_simple_timeline(slide)

    slide = add_bullet_slide(
        prs,
        title="Post-op Phase IV (12+ weeks)",
        bullets=[
            "Progressive resistance: multi-plane, endurance and power as needed",
            "Closed/open-chain integration: wall slides, push-up plus, landmine press",
            "Sport/work-specific drills when criteria met (ROM/strength/pain)",
            "Ongoing education on load management and recovery",
        ],
    )
    add_simple_timeline(slide)

    # Exercises & precautions
    add_two_content_slide(
        prs,
        title="Exercises & Precautions",
        left_bullets=[
            "Gentle pendulums, table slides, wand AAROM",
            "Scapular clocks, prone Ys/Ts (later phases)",
            "Theraband rows, ER/IR at 0–45° abduction",
            "Closed-chain weight shifts, wall push-up plus",
        ],
        right_bullets=[
            "Avoid aggressive stretching early post-op",
            "No lifting overhead or sudden jerks initially",
            "Respect pain/night pain; avoid impingement positions",
            "Follow surgeon-specific protocol and healing constraints",
        ],
    )

    # Home program & education
    add_bullet_slide(
        prs,
        title="Home Program & Education",
        bullets=[
            "Daily symptom-guided ROM (as allowed)",
            "Ice after exercises if sore; heat before if stiff",
            "Posture breaks; sleep with pillows for support",
            "Consistency over intensity; log exercises to track progress",
        ],
    )

    # Case study
    add_bullet_slide(
        prs,
        title="Case Study (Example)",
        bullets=[
            "52-year-old painter; 3 months shoulder pain, night pain",
            "Exam: painful arc 70–120°, ER weakness; +Jobe, +Hawkins",
            "MRI: medium supraspinatus tear; arthroscopic repair performed",
            "Rehab: phased protocol; at 12 wks: near-full PROM, band strengthening",
        ],
    )

    # Outcome snapshot
    add_two_content_slide(
        prs,
        title="Case Outcome (12–16 weeks)",
        left_bullets=[
            "Pain: 7/10 → 2/10 (night pain minimal)",
            "ROM: Flex 160°, Abd 150°, ER 50°",
            "Strength: ER 4/5; resumed light duties",
        ],
        right_bullets=[
            "Ongoing: progressive strengthening, overhead tolerance",
            "Return-to-work plan with graded exposure",
            "Emphasis on posture and load pacing",
        ],
    )

    # Conclusion & takeaways
    add_bullet_slide(
        prs,
        title="Conclusion & Key Takeaways",
        bullets=[
            "Early identification and load management reduce chronicity",
            "Surgery helps selected patients; rehab is essential to outcomes",
            "Phase-based PT guided by symptoms and healing timelines",
            "Communication: patient–physio–surgeon team approach",
        ],
    )

    # References
    add_bullet_slide(
        prs,
        title="References (Selected)",
        bullets=[
            "American Academy of Orthopaedic Surgeons (AAOS) guidelines",
            "Bain et al. Shoulder & Elbow texts; Kisner & Colby (Ther Ex)",
            "Recent systematic reviews on rotator cuff repair outcomes",
            "Local hospital post-op protocols (per surgeon preference)",
        ],
    )

    prs.save(output_path)


if __name__ == "__main__":
    output_file = "/workspace/Rotator_Cuff_Injury_Surgical_and_Physiotherapy_Management.pptx"
    build_presentation(output_file)
    print(f"Generated: {output_file}")

